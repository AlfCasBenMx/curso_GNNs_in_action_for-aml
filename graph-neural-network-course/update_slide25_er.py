from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# Rutas
pptx_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"
output_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"

# Cargar presentación
prs = Presentation(pptx_path)
slide_25 = prs.slides[24]

# Limpiar contenido existente (mantener solo el fondo)
shapes_to_remove = []
for shape in slide_25.shapes:
    # Verificar si es el fondo (rectángulo grande negro)
    if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE.RECTANGLE:
        if shape.width > Inches(12) and shape.height > Inches(7):
            continue  # Mantener el fondo
    shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except:
        pass

# Colores EY
AMARILLO_EY = RGBColor(255, 196, 0)
NEGRO = RGBColor(30, 30, 30)
BLANCO = RGBColor(255, 255, 255)
GRIS_OSCURO = RGBColor(50, 50, 50)
GRIS_CLARO = RGBColor(128, 128, 128)
AZUL = RGBColor(0, 120, 200)
VERDE = RGBColor(0, 180, 100)

# ============== FONDO OSCURO ==============
bg_shape = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = NEGRO
bg_shape.line.fill.background()
spTree = slide_25.shapes._spTree
sp = bg_shape._element
spTree.remove(sp)
spTree.insert(2, sp)

# ============== TÍTULO ==============
title_box = slide_25.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "SUPUESTOS Y MODELO DE DATOS"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = BLANCO

# Subtítulo
subtitle_box = slide_25.shapes.add_textbox(Inches(8.5), Inches(0.3), Inches(4), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "PoC HUNAB - Banco de Chile"
p.font.size = Pt(12)
p.font.color.rgb = AMARILLO_EY
p.alignment = PP_ALIGN.RIGHT

# ============== SECCIÓN IZQUIERDA: SUPUESTOS CLAVE ==============
# Título sección
sec_title = slide_25.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(5), Inches(0.3))
tf = sec_title.text_frame
p = tf.paragraphs[0]
p.text = "Supuestos Clave"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = AMARILLO_EY

# Lista de supuestos compacta
supuestos = [
    ("1", "Integridad Referencial", "Datos ofuscados mantienen relaciones entre tablas"),
    ("2", "Fraudes Declarados", "Incluir casos confirmados por el cliente"),
    ("3", "Resolución Dual", "Fraudes a favor cliente Y a favor banco"),
    ("4", "Formato Parquet", "Estructuras optimizadas para procesamiento ML"),
    ("5", "Etiquetado Claro", "Transacciones con flag fraude/no-fraude"),
]

y_pos = Inches(1.15)
for num, titulo, desc in supuestos:
    # Número
    num_box = slide_25.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), y_pos, Inches(0.3), Inches(0.3))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = AMARILLO_EY
    num_box.line.fill.background()
    
    num_tb = slide_25.shapes.add_textbox(Inches(0.4), y_pos, Inches(0.3), Inches(0.3))
    tf = num_tb.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NEGRO
    p.alignment = PP_ALIGN.CENTER
    
    # Título y descripción
    text_box = slide_25.shapes.add_textbox(Inches(0.8), y_pos - Inches(0.02), Inches(5.5), Inches(0.35))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{titulo}: "
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLANCO
    
    run = p.add_run()
    run.text = desc
    run.font.size = Pt(9)
    run.font.bold = False
    run.font.color.rgb = RGBColor(180, 180, 180)
    
    y_pos += Inches(0.38)

# ============== SECCIÓN DERECHA: DIAGRAMA ER ==============
# Título sección
er_title = slide_25.shapes.add_textbox(Inches(6.8), Inches(0.8), Inches(6), Inches(0.3))
tf = er_title.text_frame
p = tf.paragraphs[0]
p.text = "Modelo Entidad-Relación"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = AMARILLO_EY

# ============== ENTIDADES DEL DIAGRAMA ER ==============

def create_entity(slide, left, top, width, height, name, fields, color, is_main=False):
    """Crear una entidad del diagrama ER"""
    # Caja principal
    entity = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    entity.fill.solid()
    entity.fill.fore_color.rgb = GRIS_OSCURO
    entity.line.color.rgb = color
    entity.line.width = Pt(3 if is_main else 2)
    
    # Header de la entidad
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.35))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    
    # Nombre de entidad
    name_tb = slide.shapes.add_textbox(left, top + Inches(0.05), width, Inches(0.3))
    tf = name_tb.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NEGRO if color == AMARILLO_EY else BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Campos
    fields_tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.4), width - Inches(0.2), height - Inches(0.45))
    tf = fields_tb.text_frame
    tf.word_wrap = True
    for i, field in enumerate(fields):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # PK en negrita
        if field.startswith("PK:") or field.startswith("FK:"):
            p.text = field
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = AMARILLO_EY if field.startswith("PK:") else AZUL
        else:
            p.text = f"  {field}"
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(200, 200, 200)
        p.space_after = Pt(2)
    
    return entity

# Entidad TARJETAS (izquierda)
create_entity(slide_25, 
    Inches(6.5), Inches(1.3), Inches(2.4), Inches(1.8),
    "📇 TARJETAS", 
    ["PK: tarjeta_id", "cliente_id", "bin_tipo", "fecha_emision", "estado", "limite"],
    AMARILLO_EY, is_main=True)

# Entidad TRANSACCIONES (centro)
create_entity(slide_25,
    Inches(9.3), Inches(1.3), Inches(2.6), Inches(2.0),
    "💳 TRANSACCIONES",
    ["PK: trx_id", "FK: tarjeta_id", "fecha_hora", "monto", "comercio", "canal", "ubicacion", "dispositivo"],
    AMARILLO_EY, is_main=True)

# Entidad FRAUDES (abajo centro)
create_entity(slide_25,
    Inches(7.9), Inches(3.6), Inches(2.6), Inches(1.7),
    "🚨 FRAUDES",
    ["PK: fraude_id", "FK: trx_id", "tipo_fraude", "resolucion", "fecha_reclamo", "monto_fraude"],
    RGBColor(220, 80, 80))

# ============== LÍNEAS DE RELACIÓN ==============
def add_connector_line(slide, x1, y1, x2, y2, label="", color=GRIS_CLARO):
    """Agregar línea de conexión con etiqueta"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
        min(x1, x2), min(y1, y2), 
        Pt(3), abs(y2 - y1) if x1 == x2 else abs(x2 - x1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    
    if label:
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        lbl = slide.shapes.add_textbox(mid_x - Inches(0.3), mid_y - Inches(0.1), Inches(0.6), Inches(0.2))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(7)
        p.font.color.rgb = AMARILLO_EY
        p.alignment = PP_ALIGN.CENTER

# Línea TARJETAS -> TRANSACCIONES (horizontal)
line1 = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.9), Inches(2.1), Inches(0.4), Pt(3))
line1.fill.solid()
line1.fill.fore_color.rgb = AMARILLO_EY
line1.line.fill.background()

rel1 = slide_25.shapes.add_textbox(Inches(8.85), Inches(1.85), Inches(0.6), Inches(0.2))
tf = rel1.text_frame
p = tf.paragraphs[0]
p.text = "1:N"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = AMARILLO_EY
p.alignment = PP_ALIGN.CENTER

# Línea TRANSACCIONES -> FRAUDES (vertical hacia abajo)
line2 = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(3.3), Pt(3), Inches(0.3))
line2.fill.solid()
line2.fill.fore_color.rgb = RGBColor(220, 80, 80)
line2.line.fill.background()

# Línea horizontal conectando a FRAUDES
line2b = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(3.6), Inches(0.1), Pt(3))
line2b.fill.solid()
line2b.fill.fore_color.rgb = RGBColor(220, 80, 80)
line2b.line.fill.background()

rel2 = slide_25.shapes.add_textbox(Inches(10.55), Inches(3.35), Inches(0.6), Inches(0.2))
tf = rel2.text_frame
p = tf.paragraphs[0]
p.text = "1:0..1"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(220, 80, 80)
p.alignment = PP_ALIGN.CENTER

# ============== LEYENDA ==============
legend_box = slide_25.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.85), Inches(5.2), Inches(0.9))
legend_box.fill.solid()
legend_box.fill.fore_color.rgb = RGBColor(40, 40, 40)
legend_box.line.color.rgb = GRIS_CLARO
legend_box.line.width = Pt(1)

legend_title = slide_25.shapes.add_textbox(Inches(6.6), Inches(4.9), Inches(1), Inches(0.25))
tf = legend_title.text_frame
p = tf.paragraphs[0]
p.text = "Leyenda:"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = BLANCO

legend_items = [
    ("PK", "Primary Key (clave única)", AMARILLO_EY),
    ("FK", "Foreign Key (referencia)", AZUL),
    ("1:N", "Una tarjeta → muchas transacciones", AMARILLO_EY),
    ("1:0..1", "Una transacción → 0 ó 1 fraude", RGBColor(220, 80, 80)),
]

x_pos = Inches(6.6)
for i, (key, desc, color) in enumerate(legend_items):
    y = Inches(5.15) + (i // 2) * Inches(0.25)
    x = x_pos if i % 2 == 0 else x_pos + Inches(2.6)
    
    item_tb = slide_25.shapes.add_textbox(x, y, Inches(2.5), Inches(0.2))
    tf = item_tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{key}: "
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = color
    
    run = p.add_run()
    run.text = desc
    run.font.size = Pt(8)
    run.font.bold = False
    run.font.color.rgb = RGBColor(180, 180, 180)

# ============== TARJETAS INFERIORES ==============
tarjetas_info = [
    {
        "titulo": "📊 Datos Requeridos",
        "items": ["5 meses de historia", "BIN SIGNATURE", "29.4M transacciones"]
    },
    {
        "titulo": "🔗 Integridad",
        "items": ["IDs ofuscados pero", "consistentes entre tablas", "para análisis ML"]
    },
    {
        "titulo": "✅ Entregables",
        "items": ["Archivos Parquet", "Estructura validada", "Catálogos de datos"]
    }
]

card_width = Inches(1.95)
card_height = Inches(1.15)
card_top = Inches(5.95)
start_left = Inches(0.4)

for i, tarjeta in enumerate(tarjetas_info):
    left = start_left + i * (card_width + Inches(0.15))
    
    card_bg = slide_25.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = GRIS_OSCURO
    card_bg.line.color.rgb = AMARILLO_EY
    card_bg.line.width = Pt(1)
    
    title_tb = slide_25.shapes.add_textbox(left, card_top + Inches(0.05), card_width, Inches(0.25))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = tarjeta["titulo"]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = AMARILLO_EY
    p.alignment = PP_ALIGN.CENTER
    
    items_tb = slide_25.shapes.add_textbox(left + Inches(0.1), card_top + Inches(0.35), card_width - Inches(0.2), Inches(0.8))
    tf = items_tb.text_frame
    for j, item in enumerate(tarjeta["items"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(180, 180, 180)

# ============== NOTA AL PIE ==============
footer = slide_25.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(12), Inches(0.25))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "Validación: José Durand / Alfredo / Juvenal  |  Confirmar estructura antes del inicio de PoC"
p.font.size = Pt(8)
p.font.italic = True
p.font.color.rgb = GRIS_CLARO

# Guardar
prs.save(output_path)
print(f"✅ Slide 25 actualizada con diagrama ER")
print(f"📁 Guardado en: {output_path}")
print("\nElementos incluidos:")
print("  - Supuestos clave (lado izquierdo)")
print("  - Diagrama Entidad-Relación (lado derecho)")
print("  - Entidades: TARJETAS, TRANSACCIONES, FRAUDES, COMERCIOS")
print("  - Relaciones: 1:N, 1:1, N:1")
print("  - Leyenda explicativa")
print("  - Tarjetas de resumen")
