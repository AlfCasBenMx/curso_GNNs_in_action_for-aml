from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Ruta del archivo
pptx_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214.pptx"

# Cargar la presentación
prs = Presentation(pptx_path)

print(f"Total de diapositivas: {len(prs.slides)}")
print("="*80)

# Leer todas las diapositivas para contexto
for idx, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*80}")
    print(f"SLIDE {idx}")
    print(f"{'='*80}")
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    print(f"  - {text[:150]}{'...' if len(text) > 150 else ''}")
        
        if shape.has_table:
            print("  [TABLA]")
            table = shape.table
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text[:30])
                if row_text:
                    print(f"    | {' | '.join(row_text)} |")
