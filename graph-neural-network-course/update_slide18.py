from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from copy import deepcopy

# Ruta del archivo
pptx_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214.pptx"
output_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"

# Cargar la presentación
prs = Presentation(pptx_path)

# Obtener la slide 18 (índice 17)
slide_18 = prs.slides[17]

# Limpiar el contenido actual de la slide 18 (excepto elementos de diseño)
shapes_to_remove = []
for shape in slide_18.shapes:
    if shape.has_text_frame:
        # Verificar si es contenido que debemos limpiar
        text_content = ""
        for para in shape.text_frame.paragraphs:
            text_content += para.text
        if "Trabajar el plan de trabajo" in text_content or "Considerar Semana" in text_content or "Adriana con soporte" in text_content:
            shapes_to_remove.append(shape)

# Eliminar las formas identificadas
for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Definir el plan de trabajo
plan_trabajo = [
    {
        "semana": "Semana 0",
        "titulo": "Preparación y Kick-off",
        "actividades": [
            "Confirmación de disponibilidad de datos transaccionales",
            "Identificación y confirmación de Stakeholders (Prevención Fraude, Data & Analytics)",
            "Firma de NDA y acuerdos de confidencialidad",
            "Definición del método de recolección de logs",
            "Confirmación del BIN candidato (SIGNATURE recomendado)"
        ]
    },
    {
        "semana": "Semana 1-2",
        "titulo": "Extracción y Preparación de Datos",
        "actividades": [
            "Extracción de datos históricos (5 meses de transacciones)",
            "Validación de calidad de datos",
            "Análisis exploratorio de datos (EDA)",
            "Preparación del ambiente de desarrollo"
        ]
    },
    {
        "semana": "Semana 3-4",
        "titulo": "Desarrollo de Modelos HUNAB",
        "actividades": [
            "Entrenamiento modelo HUNAB Supervisado (XGBoost)",
            "Entrenamiento modelo HUNAB No Supervisado (KMeans)",
            "Validación cruzada y ajuste de hiperparámetros",
            "Análisis de métricas (Precisión, Recall, F1)"
        ]
    },
    {
        "semana": "Semana 5-6",
        "titulo": "Validación y Resultados",
        "actividades": [
            "Validación de resultados con equipo de Prevención de Fraudes",
            "Análisis post-mortem y detección de patrones emergentes",
            "Documentación de hallazgos y recomendaciones",
            "Presentación de resultados finales"
        ]
    }
]

# Agregar título principal
title_box = slide_18.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Plan de Trabajo - PoC HUNAB"
title_para.font.size = Pt(28)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)  # Azul oscuro

# Agregar subtítulo
subtitle_box = slide_18.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = "Metodología de Aprendizaje Continuo - Banco de Chile"
subtitle_para.font.size = Pt(14)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(128, 128, 128)

# Crear tabla de plan de trabajo
rows = len(plan_trabajo) + 1  # +1 para header
cols = 3
table_width = Inches(9.0)
table_height = Inches(4.5)

table = slide_18.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), table_width, table_height).table

# Configurar anchos de columnas
table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(5.5)

# Encabezados
headers = ["Período", "Fase", "Actividades Clave"]
header_row = table.rows[0]
for idx, header in enumerate(headers):
    cell = header_row.cells[idx]
    cell.text = header
    para = cell.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(11)
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(255, 196, 0)  # Amarillo EY

# Llenar datos del plan
for row_idx, fase in enumerate(plan_trabajo, start=1):
    row = table.rows[row_idx]
    
    # Semana
    cell_semana = row.cells[0]
    cell_semana.text = fase["semana"]
    para = cell_semana.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(10)
    para.alignment = PP_ALIGN.CENTER
    
    # Título de fase
    cell_titulo = row.cells[1]
    cell_titulo.text = fase["titulo"]
    para = cell_titulo.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(9)
    
    # Actividades
    cell_act = row.cells[2]
    tf = cell_act.text_frame
    tf.clear()
    for act_idx, actividad in enumerate(fase["actividades"]):
        if act_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {actividad}"
        p.font.size = Pt(9)
        p.space_after = Pt(2)

# Agregar nota al pie
note_box = slide_18.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.4))
note_frame = note_box.text_frame
note_para = note_frame.paragraphs[0]
note_para.text = "Duración estimada: 6 semanas | Equipo: Adriana, Alfredo, Juvenal | Entregable: Informe de resultados PoC"
note_para.font.size = Pt(10)
note_para.font.italic = True
note_para.font.color.rgb = RGBColor(100, 100, 100)

# Guardar la presentación actualizada
prs.save(output_path)
print(f"Presentación actualizada guardada en: {output_path}")
print("Slide 18 actualizada con el Plan de Trabajo completo.")
