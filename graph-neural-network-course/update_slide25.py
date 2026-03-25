from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Rutas
pptx_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"
output_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"

# Cargar presentación
prs = Presentation(pptx_path)
slide_25 = prs.slides[24]  # Índice 24 = Slide 25

# Limpiar contenido existente
shapes_to_remove = []
for shape in slide_25.shapes:
    if shape.has_text_frame:
        text_content = ""
        for para in shape.text_frame.paragraphs:
            text_content += para.text
        if any(keyword in text_content for keyword in ["Trabajar", "supuestos", "Soporte", "Alfredo", "Juvenal", "lista"]):
            shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Colores EY
AMARILLO_EY = RGBColor(255, 196, 0)
NEGRO = RGBColor(30, 30, 30)
BLANCO = RGBColor(255, 255, 255)
GRIS_OSCURO = RGBColor(50, 50, 50)
GRIS_CLARO = RGBColor(128, 128, 128)
AZUL = RGBColor(0, 120, 200)

# ============== FONDO OSCURO ==============
bg_shape = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = NEGRO
bg_shape.line.fill.background()
spTree = slide_25.shapes._spTree
sp = bg_shape._element
spTree.remove(sp)
spTree.insert(2, sp)

# ============== TÍTULO ==============
title_box = slide_25.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "SUPUESTOS Y PREMISAS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BLANCO

# Subtítulo derecha
subtitle_box = slide_25.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(4), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "PoC HUNAB - Banco de Chile"
p.font.size = Pt(14)
p.font.color.rgb = AMARILLO_EY
p.alignment = PP_ALIGN.RIGHT

# Línea decorativa
line = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(0.75), Inches(4), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = AMARILLO_EY
line.line.fill.background()

# ============== TABLA DE SUPUESTOS ==============
supuestos_data = [
    ["#", "Categoría", "Supuesto / Premisa", "Impacto"],
    ["1", "Calidad de Datos", "Los datos ofuscados mantendrán la integridad referencial para permitir el análisis de patrones", "Alto"],
    ["2", "Datos de Fraude", "Se proporcionarán los datos de fraudes declarados por el cliente (casos confirmados)", "Alto"],
    ["3", "Resolución Fraudes", "Se incluirán fraudes resueltos a favor del cliente Y a favor del banco para entrenamiento balanceado", "Alto"],
    ["4", "Formato de Datos", "Los datos serán compartidos en estructuras Parquet para optimizar procesamiento", "Medio"],
    ["5", "Volumen Histórico", "Se dispondrá de mínimo 5 meses de historia transaccional del BIN candidato", "Alto"],
    ["6", "Etiquetado ML", "Las transacciones tendrán etiquetas claras de fraude/no-fraude para el modelo supervisado", "Alto"],
    ["7", "Variables Clave", "Se incluirán variables: monto, comercio, hora, ubicación, frecuencia, dispositivo", "Alto"],
]

rows = len(supuestos_data)
cols = len(supuestos_data[0])
table_left = Inches(0.4)
table_top = Inches(1.0)
table_width = Inches(12.5)
table_height = Inches(3.2)

table_shape = slide_25.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
table = table_shape.table

# Anchos de columna
table.columns[0].width = Inches(0.5)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(8.5)
table.columns[3].width = Inches(1.5)

# Llenar tabla
for row_idx, row_data in enumerate(supuestos_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if col_idx != 2 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if row_idx == 0:  # Header
            para.font.bold = True
            para.font.size = Pt(11)
            para.font.color.rgb = NEGRO
            cell.fill.solid()
            cell.fill.fore_color.rgb = AMARILLO_EY
        else:  # Datos
            para.font.size = Pt(9)
            para.font.color.rgb = BLANCO
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRIS_OSCURO
            
            # Columna de impacto con colores
            if col_idx == 3:
                para.font.bold = True
                if cell_text == "Alto":
                    para.font.color.rgb = AMARILLO_EY
                elif cell_text == "Medio":
                    para.font.color.rgb = AZUL

# ============== TARJETAS DE CATEGORÍAS ==============
tarjetas = [
    {
        "numero": "📊",
        "titulo": "Calidad de Datos",
        "items": [
            "Integridad referencial preservada",
            "Datos ofuscados pero vinculables",
            "Sin pérdida de patrones clave"
        ],
        "color": AMARILLO_EY
    },
    {
        "numero": "🔍",
        "titulo": "Datos de Fraude",
        "items": [
            "Fraudes declarados por cliente",
            "Resolución: favor cliente/banco",
            "Etiquetado claro para ML"
        ],
        "color": AMARILLO_EY
    },
    {
        "numero": "💾",
        "titulo": "Formato y Entrega",
        "items": [
            "Estructuras Parquet",
            "5+ meses de historia",
            "BIN SIGNATURE confirmado"
        ],
        "color": AMARILLO_EY
    },
    {
        "numero": "🤖",
        "titulo": "Requisitos ML",
        "items": [
            "Variables: monto, comercio, hora",
            "Ubicación, frecuencia, dispositivo",
            "Balance fraude/no-fraude"
        ],
        "color": AMARILLO_EY
    }
]

card_width = Inches(2.9)
card_height = Inches(2.0)
card_top = Inches(4.5)
card_spacing = Inches(0.2)
start_left = Inches(0.5)

for i, tarjeta in enumerate(tarjetas):
    left = start_left + i * (card_width + card_spacing)
    
    # Fondo de tarjeta
    card_bg = slide_25.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = GRIS_OSCURO
    card_bg.line.color.rgb = tarjeta["color"]
    card_bg.line.width = Pt(2)
    
    # Título de tarjeta
    title_tb = slide_25.shapes.add_textbox(left, card_top + Inches(0.1), card_width, Inches(0.3))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = tarjeta["titulo"]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = tarjeta["color"]
    p.alignment = PP_ALIGN.CENTER
    
    # Línea separadora
    sep_line = slide_25.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.3), card_top + Inches(0.4), card_width - Inches(0.6), Pt(2))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = tarjeta["color"]
    sep_line.line.fill.background()
    
    # Items
    items_tb = slide_25.shapes.add_textbox(left + Inches(0.15), card_top + Inches(0.5), card_width - Inches(0.3), Inches(1.4))
    tf = items_tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(tarjeta["items"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.space_after = Pt(4)

# ============== NOTA AL PIE ==============
footer = slide_25.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.3))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "Validación de supuestos: José Durand / Alfredo Castañeda / Juvenal  |  Estos supuestos deben ser confirmados antes del inicio de la PoC"
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = GRIS_CLARO
p.alignment = PP_ALIGN.CENTER

# Guardar
prs.save(output_path)
print(f"✅ Slide 25 actualizada con supuestos y premisas")
print(f"📁 Guardado en: {output_path}")
print("\nSupuestos incluidos basados en conversación con José Durand:")
print("  - Integridad referencial de datos ofuscados")
print("  - Fraudes declarados por cliente")
print("  - Resolución a favor cliente Y banco")
print("  - Estructuras Parquet para compartir datos")
print("  - Variables clave para ML")
print("  - Etiquetado para modelo supervisado")
