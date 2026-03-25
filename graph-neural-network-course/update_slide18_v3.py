from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Rutas
pptx_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214.pptx"
output_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"

# Cargar presentación
prs = Presentation(pptx_path)
slide_18 = prs.slides[17]

# Limpiar contenido existente
shapes_to_remove = []
for shape in slide_18.shapes:
    if shape.has_text_frame:
        text_content = ""
        for para in shape.text_frame.paragraphs:
            text_content += para.text
        if "Trabajar el plan" in text_content or "Considerar Semana" in text_content or "Adriana con soporte" in text_content:
            shapes_to_remove.append(shape)
    if shape.has_table:
        shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Colores EY
AMARILLO_EY = RGBColor(255, 196, 0)
NEGRO = RGBColor(30, 30, 30)
BLANCO = RGBColor(255, 255, 255)
GRIS_OSCURO = RGBColor(50, 50, 50)
GRIS_CLARO = RGBColor(128, 128, 128)

# ============== FONDO OSCURO PARA LA SLIDE ==============
bg_shape = slide_18.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = NEGRO
bg_shape.line.fill.background()
# Mover al fondo
spTree = slide_18.shapes._spTree
sp = bg_shape._element
spTree.remove(sp)
spTree.insert(2, sp)

# ============== TÍTULO ==============
title_box = slide_18.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "PLAN DE TRABAJO"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BLANCO

# Subtítulo derecha
subtitle_box = slide_18.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(4), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "PoC HUNAB - Banco de Chile"
p.font.size = Pt(14)
p.font.color.rgb = AMARILLO_EY
p.alignment = PP_ALIGN.RIGHT

# Línea decorativa amarilla
line = slide_18.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(0.75), Inches(4), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = AMARILLO_EY
line.line.fill.background()

# ============== TABLA PLAN DE TRABAJO ==============
# Datos del plan
plan_data = [
    ["Semana", "Fase", "Actividades Clave", "Entregables", "Responsables"],
    ["Semana 0", "Kick-off", "Confirmación datos y stakeholders\nFirma NDA, definición método logs", "Acta de inicio", "Marlene / José"],
    ["Semana 1-2", "Datos", "Extracción histórica 5 meses\nValidación calidad, EDA", "Dataset limpio", "Alfredo / Luis"],
    ["Semana 3-4", "Modelos", "Entrenamiento XGBoost y KMeans\nAjuste hiperparámetros", "Modelos entrenados", "Alfredo / Angelo"],
    ["Semana 5-6", "Validación", "Análisis post-mortem\nPresentación resultados", "Informe final PoC", "Equipo completo"],
]

rows = len(plan_data)
cols = len(plan_data[0])
table_left = Inches(0.4)
table_top = Inches(1.1)
table_width = Inches(12.5)
table_height = Inches(2.8)

table_shape = slide_18.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
table = table_shape.table

# Anchos de columna
table.columns[0].width = Inches(1.4)
table.columns[1].width = Inches(1.4)
table.columns[2].width = Inches(4.5)
table.columns[3].width = Inches(2.5)
table.columns[4].width = Inches(2.7)

# Llenar tabla
for row_idx, row_data in enumerate(plan_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        
        # Formato de celda
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if row_idx == 0:  # Header
            para.font.bold = True
            para.font.size = Pt(11)
            para.font.color.rgb = NEGRO
            cell.fill.solid()
            cell.fill.fore_color.rgb = AMARILLO_EY
        else:  # Datos
            para.font.size = Pt(9)
            para.font.color.rgb = BLANCO
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRIS_OSCURO
            
            # Para celdas con múltiples líneas
            if "\n" in cell_text:
                cell.text = ""
                tf = cell.text_frame
                lines = cell_text.split("\n")
                for i, line_text in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = "• " + line_text
                    p.font.size = Pt(9)
                    p.font.color.rgb = BLANCO
                    p.alignment = PP_ALIGN.LEFT

# ============== ICONOS DE TEXTO LATERAL ==============
icon_label = slide_18.shapes.add_textbox(Inches(0.1), Inches(1.8), Inches(0.3), Inches(2))
tf = icon_label.text_frame
p = tf.paragraphs[0]
p.text = "Cronograma\nde Ejecución"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = AMARILLO_EY
# Rotar texto (simulado con posición)

# ============== TARJETAS INFERIORES ==============
tarjetas = [
    {
        "icono": "📊",
        "titulo": "Semana 0",
        "subtitulo": "Preparación",
        "descripcion": "Confirmar datos listos\ny stakeholders\ndisponibles"
    },
    {
        "icono": "🔍",
        "titulo": "Semana 1-2", 
        "subtitulo": "Exploración",
        "descripcion": "Análisis de 5M+\ntransacciones\ny 200K tarjetas"
    },
    {
        "icono": "🤖",
        "titulo": "Semana 3-4",
        "subtitulo": "Modelado",
        "descripcion": "Supervisado: 90% Precisión\nNo Supervisado:\ndetección anomalías"
    },
    {
        "icono": "✅",
        "titulo": "Semana 5-6",
        "subtitulo": "Resultados",
        "descripcion": "Ahorro potencial:\n$80K USD/mes\n686 fraudes evitados"
    }
]

card_width = Inches(2.9)
card_height = Inches(1.8)
card_top = Inches(4.2)
card_spacing = Inches(0.2)
start_left = Inches(0.5)

for i, tarjeta in enumerate(tarjetas):
    left = start_left + i * (card_width + card_spacing)
    
    # Fondo de tarjeta (rectángulo redondeado)
    card_bg = slide_18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = GRIS_OSCURO
    card_bg.line.color.rgb = GRIS_CLARO
    card_bg.line.width = Pt(1)
    
    # Círculo con icono
    circle_size = Inches(0.6)
    circle = slide_18.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.15), card_top + Inches(0.15), circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(80, 80, 80)
    circle.line.color.rgb = AMARILLO_EY
    circle.line.width = Pt(2)
    
    # Título de tarjeta
    title_tb = slide_18.shapes.add_textbox(left, card_top + Inches(0.8), card_width, Inches(0.25))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = tarjeta["titulo"]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = AMARILLO_EY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    sub_tb = slide_18.shapes.add_textbox(left, card_top + Inches(1.0), card_width, Inches(0.2))
    tf = sub_tb.text_frame
    p = tf.paragraphs[0]
    p.text = tarjeta["subtitulo"]
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Descripción
    desc_tb = slide_18.shapes.add_textbox(left + Inches(0.1), card_top + Inches(1.2), card_width - Inches(0.2), Inches(0.6))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(tarjeta["descripcion"].split("\n")):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

# ============== NOTA AL PIE ==============
footer = slide_18.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.3))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "Duración total: 6 semanas  |  Equipo EY: Marlene, José, Oscar, Alfredo, Luis  |  Entregable final: Informe PoC con hallazgos y recomendaciones"
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = GRIS_CLARO
p.alignment = PP_ALIGN.CENTER

# Guardar
prs.save(output_path)
print(f"✅ Presentación actualizada guardada en:\n{output_path}")
print("\nSlide 18 actualizada con formato visual mejorado:")
print("  - Fondo oscuro")
print("  - Tabla con headers amarillos")
print("  - 4 tarjetas informativas con métricas clave")
