from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Rutas
pptx_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"
output_path = r"C:\Users\RX971WG\OneDrive - EY\Documents\graph-neural-network-course\Banco de Chile - EY - Casos de Uso IA - Fraude 20260214_updated.pptx"

# Cargar presentación
prs = Presentation(pptx_path)
slide_23 = prs.slides[22]  # Índice 22 = Slide 23

# Limpiar contenido existente
shapes_to_remove = []
for shape in slide_23.shapes:
    if shape.has_text_frame:
        text_content = ""
        for para in shape.text_frame.paragraphs:
            text_content += para.text
        # Remover contenido de próximos pasos
        if any(keyword in text_content for keyword in ["Hunab", "PoC", "Próximos", "NDA", "Confirmar", "Definir", "Identificar", "Revisión", "stakeholders", "BIN"]):
            shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Colores EY
AMARILLO_EY = RGBColor(255, 196, 0)
NEGRO = RGBColor(30, 30, 30)
BLANCO = RGBColor(255, 255, 255)
GRIS_OSCURO = RGBColor(50, 50, 50)
GRIS_CLARO = RGBColor(128, 128, 128)
VERDE = RGBColor(0, 180, 100)

# ============== FONDO OSCURO ==============
bg_shape = slide_23.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = NEGRO
bg_shape.line.fill.background()
spTree = slide_23.shapes._spTree
sp = bg_shape._element
spTree.remove(sp)
spTree.insert(2, sp)

# ============== TÍTULO ==============
title_box = slide_23.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "PRÓXIMOS PASOS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BLANCO

# Subtítulo derecha
subtitle_box = slide_23.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(4), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "Hunab – PoC – Banco de Chile"
p.font.size = Pt(14)
p.font.color.rgb = AMARILLO_EY
p.alignment = PP_ALIGN.RIGHT

# Línea decorativa
line = slide_23.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(0.75), Inches(4), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = AMARILLO_EY
line.line.fill.background()

# ============== TABLA DE PRÓXIMOS PASOS ==============
pasos_data = [
    ["#", "Acción", "Descripción", "Responsable", "Estado"],
    ["1", "Firma NDA", "Avanzar con la firma del acuerdo de confidencialidad", "Legal / Marlene", "Pendiente"],
    ["2", "Modalidad Ejecución", "Confirmar modalidad de ejecución recomendada", "José Durand", "Pendiente"],
    ["3", "Método Logs", "Definir método de recolección de logs para la PoC", "Equipo EY", "Pendiente"],
    ["4", "Stakeholders", "Identificar stakeholders de Prevención y Data & Analytics", "Banco de Chile", "Pendiente"],
    ["5", "BIN Candidato", "Confirmar BIN SIGNATURE como candidato para PoC", "Alfredo / Juvenal", "Pendiente"],
    ["6", "Fecha Inicio", "Confirmar fecha tentativa de inicio de PoC", "Ambos equipos", "Pendiente"],
]

rows = len(pasos_data)
cols = len(pasos_data[0])
table_left = Inches(0.4)
table_top = Inches(1.1)
table_width = Inches(12.5)
table_height = Inches(3.0)

table_shape = slide_23.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
table = table_shape.table

# Anchos de columna
table.columns[0].width = Inches(0.5)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(5.5)
table.columns[3].width = Inches(2.2)
table.columns[4].width = Inches(2.3)

# Llenar tabla
for row_idx, row_data in enumerate(pasos_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if row_idx == 0:  # Header
            para.font.bold = True
            para.font.size = Pt(11)
            para.font.color.rgb = NEGRO
            cell.fill.solid()
            cell.fill.fore_color.rgb = AMARILLO_EY
        else:  # Datos
            para.font.size = Pt(10)
            para.font.color.rgb = BLANCO
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRIS_OSCURO
            
            # Columna de estado con color especial
            if col_idx == 4:
                para.font.color.rgb = AMARILLO_EY
                para.font.italic = True

# ============== TARJETAS INFERIORES ==============
tarjetas = [
    {
        "numero": "01",
        "titulo": "Acuerdos Legales",
        "descripcion": "NDA firmado\npara proteger\ninformación sensible",
        "color": AMARILLO_EY
    },
    {
        "numero": "02",
        "titulo": "Datos Listos",
        "descripcion": "5 meses históricos\nBIN SIGNATURE\n29.4M transacciones",
        "color": AMARILLO_EY
    },
    {
        "numero": "03",
        "titulo": "Equipo Confirmado",
        "descripcion": "Stakeholders de\nPrevención Fraude\ny Data & Analytics",
        "color": AMARILLO_EY
    },
    {
        "numero": "04",
        "titulo": "Inicio PoC",
        "descripcion": "Fecha acordada\nAmbiente preparado\nKick-off ejecutado",
        "color": VERDE
    }
]

card_width = Inches(2.9)
card_height = Inches(1.9)
card_top = Inches(4.4)
card_spacing = Inches(0.2)
start_left = Inches(0.5)

for i, tarjeta in enumerate(tarjetas):
    left = start_left + i * (card_width + card_spacing)
    
    # Fondo de tarjeta
    card_bg = slide_23.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = GRIS_OSCURO
    card_bg.line.color.rgb = tarjeta["color"]
    card_bg.line.width = Pt(2)
    
    # Número grande
    num_tb = slide_23.shapes.add_textbox(left + Inches(0.15), card_top + Inches(0.1), Inches(0.6), Inches(0.5))
    tf = num_tb.text_frame
    p = tf.paragraphs[0]
    p.text = tarjeta["numero"]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = tarjeta["color"]
    
    # Círculo decorativo
    circle = slide_23.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(2.0), card_top + Inches(0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(80, 80, 80)
    circle.line.color.rgb = tarjeta["color"]
    circle.line.width = Pt(2)
    
    # Título de tarjeta
    title_tb = slide_23.shapes.add_textbox(left, card_top + Inches(0.7), card_width, Inches(0.3))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = tarjeta["titulo"]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = tarjeta["color"]
    p.alignment = PP_ALIGN.CENTER
    
    # Descripción
    desc_tb = slide_23.shapes.add_textbox(left + Inches(0.1), card_top + Inches(1.0), card_width - Inches(0.2), Inches(0.8))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    for j, line_text in enumerate(tarjeta["descripcion"].split("\n")):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

# ============== FLECHA DE PROGRESO ==============
# Añadir flechas entre tarjetas
for i in range(3):
    left = start_left + (i + 1) * (card_width + card_spacing) - Inches(0.15)
    arrow = slide_23.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left - Inches(0.1), card_top + Inches(0.85), Inches(0.25), Inches(0.2))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = AMARILLO_EY
    arrow.line.fill.background()

# ============== NOTA AL PIE ==============
footer = slide_23.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.3))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "Objetivo: Completar todos los pasos previos antes del inicio de la Semana 0  |  Revisión y confirmación: Alfredo / Juvenal"
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = GRIS_CLARO
p.alignment = PP_ALIGN.CENTER

# Guardar
prs.save(output_path)
print(f"✅ Slide 23 actualizada con formato visual mejorado")
print(f"📁 Guardado en: {output_path}")
print("\nElementos añadidos:")
print("  - Fondo oscuro")
print("  - Tabla de próximos pasos con headers amarillos")
print("  - 4 tarjetas de hitos clave con flechas de progreso")
